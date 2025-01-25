from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional, Dict
from user_research_platform import UserResearchPlatform
import os
from dotenv import load_dotenv
from fastapi.responses import JSONResponse
from document_processor import DocumentProcessor  # Updated import
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import markdown
import io  # Add this import

load_dotenv()

app = FastAPI()
document_processor = DocumentProcessor()  # Add this line


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


class ProjectInfo(BaseModel):
    project_name: str
    goal: str
    target_audience: str
    api_key: Optional[str] = None  # Add this line



class ResponseModel(BaseModel):
    response: str


research_sessions = {}


def extract_text_from_document(file_content: bytes, file_type: str) -> str:
    """Extract text from various document formats"""
    print(f"Processing file type: {file_type}")  # Debug log

    # Normalize file type to lowercase and remove any leading dots
    file_type = file_type.lower().strip('.')

    try:
        if file_type.endswith('pdf'):
            print("Processing as PDF")  # Debug log
            pdf_reader = PdfReader(io.BytesIO(file_content))
            text_content = ""
            for page in pdf_reader.pages:
                text_content += page.extract_text() + "\n"
            return text_content

        elif file_type.endswith('docx'):
            print("Processing as DOCX")  # Debug log
            doc = Document(io.BytesIO(file_content))
            text_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text_content

        elif file_type.endswith('pptx'):
            print("Processing as PPTX")  # Debug log
            prs = Presentation(io.BytesIO(file_content))
            text_content = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
            return text_content

        elif file_type.endswith('md'):
            print("Processing as Markdown")  # Debug log
            text_content = file_content.decode('utf-8')
            html = markdown.markdown(text_content)
            from bs4 import BeautifulSoup
            text = BeautifulSoup(html, 'html.parser').get_text()
            return text

        elif file_type.endswith('txt'):
            print("Processing as TXT")  # Debug log
            return file_content.decode('utf-8')

        else:
            supported_formats = ['.pdf', '.docx', '.pptx', '.md', '.txt']
            raise ValueError(
                f"Unsupported file type: {file_type}. Supported formats are: {', '.join(supported_formats)}")

    except Exception as e:
        print(f"Error processing file: {str(e)}")  # Debug log
        raise ValueError(f"Error processing {file_type} file: {str(e)}")


@app.post("/api/start-project")
async def start_project(
        project_name: str = Form(...),
        target_audience: str = Form(...),
        goal: str = Form(...),
        improvement_objective: Optional[str] = Form(None),
        discovery_type: Optional[str] = Form(None),
        domain: Optional[str] = Form(None),
        specific_goal: Optional[str] = Form(None),
        product_doc: Optional[UploadFile] = None,
        product_name: Optional[str] = Form(None),
        product_context: Optional[str] = Form(None)
):
    try:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise HTTPException(status_code=500, detail="OpenAI API key not found")

        platform = UserResearchPlatform(api_key)
        platform.project_info = {
                     project_name: str = Form(...),
                    target_audience: str = Form(...),
                    goal: str = Form(...),
                    improvement_objective: Optional[str] = Form(None),
                    discovery_type: Optional[str] = Form(None),
                    domain: Optional[str] = Form(None),
                    specific_goal: Optional[str] = Form(None),
                    product_doc: Optional[UploadFile] = None,
                    product_name: Optional[str] = Form(None),
                    product_context: Optional[str] = Form(None)
        }

        # Process the document if provided
        if product_doc:
            print(f"Received file: {product_doc.filename}")  # Debug log
            contents = await product_doc.read()
            try:
                # Get file extension
                file_extension = product_doc.filename.split('.')[-1] if '.' in product_doc.filename else ''
                if not file_extension:
                    raise ValueError("File has no extension")

                # Extract text from the document
                text_content = extract_text_from_document(contents, file_extension)

                if not text_content.strip():
                    raise ValueError("No text content could be extracted from the document")

                # Set product context with the extracted text
                success = platform.set_product_context(text_content)
                if not success:
                    raise ValueError("Failed to process document context")

            except ValueError as e:
                raise HTTPException(
                    status_code=400,
                    detail=f"Document processing error: {str(e)}"
                )
            except Exception as e:
                raise HTTPException(
                    status_code=500,
                    detail=f"Server error processing document: {str(e)}"
                )

            if goal == "diagnostic" and product_context:
            platform.set_product_context(product_context)

        question = platform.generate_next_question()
        if not question:
            raise HTTPException(status_code=500, detail="Failed to generate question")

        session_id = f"{project_name.lower().replace(' ', '_')}_{len(research_sessions)}"
        research_sessions[session_id] = platform

        return {
            "session_id": session_id,
            "question": question,
            "api_key": api_key
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# Add new endpoint for document upload
@app.post("/api/upload-document")
async def upload_document(file: UploadFile = File(...)):
    try:
        content = await file.read()
        content_str = content.decode()

        # Process the document based on its content type
        result = document_processor.process_document(
            content_str,
            file.content_type
        )

        if "error" in result:
            raise HTTPException(status_code=400, detail=result["error"])

        return {
            "message": "Document processed successfully",
            "result": result
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# main.py
@app.post("/api/submit-response/{session_id}")
async def submit_response(session_id: str, response: ResponseModel):
    try:
        platform = research_sessions.get(session_id)
        if not platform:
            raise HTTPException(status_code=404, detail="Session not found")

        platform.conversation_history.append({
            "question": platform.current_question,
            "response": response.response
        })

        next_question = platform.generate_next_question()

        if next_question == "RESCHEDULE":
            return {
                "status": "reschedule",
                "question": "I notice we might not be getting detailed responses. Would you prefer to continue this conversation at a better time?"
            }

        if next_question == "END_INTERVIEW":
            analysis = platform.analyze_interview(platform.conversation_history)
            platform.save_results(analysis)
            del research_sessions[session_id]
            return {
                "status": "ended",
                "analysis": analysis,
                "message": "Interview ended"
            }

        return {
            "status": "continue",
            "question": next_question,
            "can_finish": len(platform.conversation_history) >= 2
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/analyze/{session_id}")
async def analyze_interview(session_id: str):
    try:
        platform = research_sessions.get(session_id)
        if not platform:
            return JSONResponse(
                status_code=404,
                content={"detail": "Session expired. Please start a new interview."}
            )

        # Only analyze if we have responses
        if len(platform.conversation_history) < 1:
            return JSONResponse(
                status_code=400,
                content={"detail": "Not enough responses to analyze."}
            )



        analysis = platform.analyze_interview(platform.conversation_history)
        filepath = document_processor.save_results(analysis)
        print(f"Analysis saved to: {filepath}")

      ### platform.save_results(analysis)

        # Save analysis results using document processor


        # Don't delete the session immediately
        # research_sessions[session_id] = platform

        return analysis
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"detail": f"Analysis failed: {str(e)}"}
        )


if __name__ == "__main__":
    import uvicorn

    print("Starting server...")
    print("Make sure your .env file contains OPENAI_API_KEY")
    uvicorn.run(app, host="0.0.0.0", port=8000)
